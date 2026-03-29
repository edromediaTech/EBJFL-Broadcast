"""Media Hub - Upload, conversion et projection de fichiers médias.

Gère l'upload multi-device, la conversion PPTX/PDF en images,
et la navigation slide par slide via WebSocket.
"""

import os
import json
import uuid
import shutil
import asyncio
import threading
from dataclasses import dataclass, asdict, field
from pathlib import Path
from datetime import datetime

UPLOAD_DIR = Path("assets/uploads")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

ALLOWED_EXTENSIONS = {
    "image": {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"},
    "video": {".mp4", ".webm", ".avi", ".mov", ".mkv"},
    "presentation": {".pptx", ".ppt"},
    "document": {".pdf", ".txt"},
}

ALL_ALLOWED = set()
for exts in ALLOWED_EXTENSIONS.values():
    ALL_ALLOWED.update(exts)

MAX_FILE_SIZE = 200 * 1024 * 1024  # 200 MB


@dataclass
class MediaFile:
    id: str
    filename: str
    original_name: str
    file_type: str          # image, video, presentation, document
    extension: str
    size: int
    path: str
    thumbnail: str = ""
    slides: list[str] = field(default_factory=list)
    total_slides: int = 0
    uploaded_at: str = ""
    status: str = "ready"   # uploading, converting, ready, error
    error: str = ""


@dataclass
class SlideshowState:
    active: bool = False
    media_id: str = ""
    current_slide: int = 0
    total_slides: int = 0
    slide_url: str = ""
    title: str = ""


class MediaHub:
    """Gère les fichiers médias uploadés."""

    def __init__(self):
        self.files: dict[str, MediaFile] = {}
        self.slideshow = SlideshowState()
        self._listeners: list = []
        self._load_existing()

    def _load_existing(self):
        """Charge les fichiers existants au démarrage."""
        if not UPLOAD_DIR.exists():
            return
        manifest = UPLOAD_DIR / "manifest.json"
        if manifest.exists():
            try:
                with open(manifest, "r", encoding="utf-8") as f:
                    data = json.load(f)
                for item in data:
                    mf = MediaFile(**item)
                    self.files[mf.id] = mf
            except Exception:
                pass

    def _save_manifest(self):
        manifest = UPLOAD_DIR / "manifest.json"
        data = [asdict(mf) for mf in self.files.values()]
        with open(manifest, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def add_listener(self, callback):
        self._listeners.append(callback)

    async def _notify(self, msg_type: str, data: dict):
        payload = {"type": msg_type, **data}
        for cb in list(self._listeners):
            try:
                await cb(payload)
            except Exception:
                self._listeners.remove(cb)

    def get_file_type(self, ext: str) -> str:
        ext = ext.lower()
        for ftype, exts in ALLOWED_EXTENSIONS.items():
            if ext in exts:
                return ftype
        return "unknown"

    async def save_upload(self, filename: str, content: bytes) -> MediaFile:
        """Sauvegarde un fichier uploadé."""
        ext = Path(filename).suffix.lower()
        if ext not in ALL_ALLOWED:
            raise ValueError(f"Type de fichier non supporté: {ext}")
        if len(content) > MAX_FILE_SIZE:
            raise ValueError(f"Fichier trop volumineux (max {MAX_FILE_SIZE // 1024 // 1024} MB)")

        file_id = str(uuid.uuid4())[:8]
        date_dir = datetime.now().strftime("%Y-%m-%d")
        save_dir = UPLOAD_DIR / date_dir
        save_dir.mkdir(parents=True, exist_ok=True)

        safe_name = f"{file_id}_{Path(filename).stem}{ext}"
        file_path = save_dir / safe_name

        with open(file_path, "wb") as f:
            f.write(content)

        ftype = self.get_file_type(ext)
        mf = MediaFile(
            id=file_id,
            filename=safe_name,
            original_name=filename,
            file_type=ftype,
            extension=ext,
            size=len(content),
            path=str(file_path.relative_to(Path("."))).replace("\\", "/"),
            uploaded_at=datetime.now().isoformat(),
            status="ready",
        )

        # Generate thumbnail for images
        if ftype == "image":
            mf.thumbnail = self._make_thumbnail(file_path, save_dir, file_id)

        # Convert presentations/documents
        if ftype == "presentation":
            mf.status = "converting"
            self.files[file_id] = mf
            self._save_manifest()
            threading.Thread(target=self._convert_pptx, args=(mf,), daemon=True).start()
            return mf

        if ftype == "document" and ext == ".pdf":
            mf.status = "converting"
            self.files[file_id] = mf
            self._save_manifest()
            threading.Thread(target=self._convert_pdf, args=(mf,), daemon=True).start()
            return mf

        self.files[file_id] = mf
        self._save_manifest()
        return mf

    def _make_thumbnail(self, src: Path, dest_dir: Path, file_id: str) -> str:
        """Crée une thumbnail 320x180."""
        try:
            from PIL import Image
            img = Image.open(src)
            img.thumbnail((320, 180))
            thumb_path = dest_dir / f"{file_id}_thumb.jpg"
            img.convert("RGB").save(thumb_path, "JPEG", quality=80)
            return str(thumb_path.relative_to(Path("."))).replace("\\", "/")
        except Exception:
            return ""

    def _convert_pptx(self, mf: MediaFile):
        """Convertit un PPTX en images PNG."""
        try:
            from pptx import Presentation
            from PIL import Image

            src = Path(mf.path)
            slides_dir = src.parent / f"{mf.id}_slides"
            slides_dir.mkdir(exist_ok=True)

            # Try comtypes (Windows PowerPoint COM) for high quality
            slides = self._convert_pptx_com(src, slides_dir)

            if not slides:
                # Fallback: extract slide content as simple images
                slides = self._convert_pptx_basic(src, slides_dir)

            mf.slides = [str(Path(s).relative_to(Path("."))).replace("\\", "/") for s in slides]
            mf.total_slides = len(slides)
            mf.status = "ready"

            if slides:
                mf.thumbnail = self._make_thumbnail(Path(slides[0]), slides_dir.parent, mf.id)

            self._save_manifest()

        except Exception as e:
            mf.status = "error"
            mf.error = str(e)
            self._save_manifest()

    def _convert_pptx_com(self, src: Path, dest_dir: Path) -> list[str]:
        """Conversion via PowerPoint COM (Windows, haute qualité)."""
        try:
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1

            abs_src = str(src.resolve())
            presentation = powerpoint.Presentations.Open(abs_src, WithWindow=False)

            slides = []
            for i, slide in enumerate(presentation.Slides):
                out = str((dest_dir / f"slide_{i+1:03d}.png").resolve())
                slide.Export(out, "PNG", 1920, 1080)
                slides.append(out)

            presentation.Close()
            powerpoint.Quit()
            return slides
        except Exception:
            return []

    def _convert_pptx_basic(self, src: Path, dest_dir: Path) -> list[str]:
        """Conversion basique via python-pptx + Pillow (fallback)."""
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw, ImageFont

            prs = Presentation(str(src))
            slides = []

            for i, slide in enumerate(prs.slides):
                # Create a simple rendered slide
                img = Image.new("RGB", (1920, 1080), "#1a1a2e")
                draw = ImageDraw.Draw(img)

                y_pos = 100
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                try:
                                    font = ImageFont.truetype("segoeui.ttf", 48)
                                except Exception:
                                    font = ImageFont.load_default()
                                draw.text((100, y_pos), text, fill="#ffffff", font=font)
                                y_pos += 70

                out = str(dest_dir / f"slide_{i+1:03d}.png")
                img.save(out)
                slides.append(out)

            return slides
        except Exception:
            return []

    def _convert_pdf(self, mf: MediaFile):
        """Convertit un PDF en images PNG."""
        try:
            src = Path(mf.path)
            slides_dir = src.parent / f"{mf.id}_slides"
            slides_dir.mkdir(exist_ok=True)

            slides = []
            try:
                # Try pdf2image if available
                from pdf2image import convert_from_path
                images = convert_from_path(str(src), dpi=150, size=(1920, 1080))
                for i, img in enumerate(images):
                    out = str(slides_dir / f"page_{i+1:03d}.png")
                    img.save(out, "PNG")
                    slides.append(out)
            except ImportError:
                # Fallback: just note it's a PDF
                from PIL import Image, ImageDraw, ImageFont
                img = Image.new("RGB", (1920, 1080), "#1a1a2e")
                draw = ImageDraw.Draw(img)
                try:
                    font = ImageFont.truetype("segoeui.ttf", 36)
                except Exception:
                    font = ImageFont.load_default()
                draw.text((200, 400), f"PDF: {mf.original_name}", fill="#ffffff", font=font)
                draw.text((200, 500), "Conversion PDF nécessite pdf2image + poppler", fill="#a6adc8", font=font)
                out = str(slides_dir / "page_001.png")
                img.save(out)
                slides.append(out)

            mf.slides = [str(Path(s).relative_to(Path("."))).replace("\\", "/") for s in slides]
            mf.total_slides = len(slides)
            mf.status = "ready"
            if slides:
                mf.thumbnail = self._make_thumbnail(Path(slides[0]), slides_dir.parent, mf.id)
            self._save_manifest()

        except Exception as e:
            mf.status = "error"
            mf.error = str(e)
            self._save_manifest()

    def get_file(self, file_id: str) -> MediaFile | None:
        return self.files.get(file_id)

    def list_files(self, file_type: str = "") -> list[dict]:
        files = sorted(self.files.values(), key=lambda f: f.uploaded_at, reverse=True)
        if file_type:
            files = [f for f in files if f.file_type == file_type]
        return [asdict(f) for f in files]

    def delete_file(self, file_id: str) -> bool:
        mf = self.files.pop(file_id, None)
        if not mf:
            return False
        # Delete physical files
        try:
            fpath = Path(mf.path)
            if fpath.exists():
                fpath.unlink()
            if mf.thumbnail:
                tp = Path(mf.thumbnail)
                if tp.exists():
                    tp.unlink()
            # Delete slides dir
            slides_dir = fpath.parent / f"{mf.id}_slides"
            if slides_dir.exists():
                shutil.rmtree(slides_dir)
        except Exception:
            pass
        self._save_manifest()
        return True

    # ── Slideshow Control ──

    async def project_file(self, file_id: str, slide_index: int = 0):
        mf = self.files.get(file_id)
        if not mf:
            return

        if mf.file_type == "image":
            self.slideshow = SlideshowState(
                active=True, media_id=file_id,
                current_slide=0, total_slides=1,
                slide_url=f"/{mf.path}", title=mf.original_name,
            )
        elif mf.slides:
            idx = max(0, min(slide_index, len(mf.slides) - 1))
            self.slideshow = SlideshowState(
                active=True, media_id=file_id,
                current_slide=idx, total_slides=len(mf.slides),
                slide_url=f"/{mf.slides[idx]}", title=mf.original_name,
            )
        elif mf.file_type == "video":
            self.slideshow = SlideshowState(
                active=True, media_id=file_id,
                current_slide=0, total_slides=0,
                slide_url=f"/{mf.path}", title=mf.original_name,
            )

        await self._notify("slideshow", asdict(self.slideshow))

    async def slide_next(self):
        if not self.slideshow.active:
            return
        mf = self.files.get(self.slideshow.media_id)
        if not mf or not mf.slides:
            return
        idx = min(self.slideshow.current_slide + 1, len(mf.slides) - 1)
        self.slideshow.current_slide = idx
        self.slideshow.slide_url = f"/{mf.slides[idx]}"
        await self._notify("slideshow", asdict(self.slideshow))

    async def slide_prev(self):
        if not self.slideshow.active:
            return
        mf = self.files.get(self.slideshow.media_id)
        if not mf or not mf.slides:
            return
        idx = max(self.slideshow.current_slide - 1, 0)
        self.slideshow.current_slide = idx
        self.slideshow.slide_url = f"/{mf.slides[idx]}"
        await self._notify("slideshow", asdict(self.slideshow))

    async def slide_goto(self, index: int):
        if not self.slideshow.active:
            return
        mf = self.files.get(self.slideshow.media_id)
        if not mf or not mf.slides:
            return
        idx = max(0, min(index, len(mf.slides) - 1))
        self.slideshow.current_slide = idx
        self.slideshow.slide_url = f"/{mf.slides[idx]}"
        await self._notify("slideshow", asdict(self.slideshow))

    async def stop_slideshow(self):
        self.slideshow = SlideshowState()
        await self._notify("slideshow", asdict(self.slideshow))

    def get_slideshow_state(self) -> dict:
        return asdict(self.slideshow)


# Singleton
media_hub = MediaHub()
